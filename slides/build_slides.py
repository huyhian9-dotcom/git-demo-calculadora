"""Gera a apresentação PPTX da aula de versionamento.

Uso:
    python slides/build_slides.py

Gera o arquivo slides/apresentacao.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


COR_FUNDO = RGBColor(0x1E, 0x1E, 0x2E)
COR_ACENTO = RGBColor(0xF9, 0xA0, 0x26)
COR_TEXTO = RGBColor(0xE4, 0xE4, 0xE7)
COR_TITULO = RGBColor(0xFF, 0xFF, 0xFF)
COR_CODIGO_BG = RGBColor(0x11, 0x11, 0x1B)
COR_CODIGO_FG = RGBColor(0x10, 0xD9, 0x90)
COR_VERMELHO = RGBColor(0xEF, 0x44, 0x44)
COR_VERDE = RGBColor(0x10, 0xD9, 0x90)


def pintar_fundo(slide, cor=COR_FUNDO):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = cor
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_faixa_acento(slide, topo=True):
    y = 0 if topo else Inches(7.3)
    faixa = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, Inches(13.333), Inches(0.2))
    faixa.fill.solid()
    faixa.fill.fore_color.rgb = COR_ACENTO
    faixa.line.fill.background()


def add_texto(slide, texto, left, top, width, height, tamanho=18, bold=False,
              cor=COR_TEXTO, alinhar=PP_ALIGN.LEFT, fonte="Calibri"):
    caixa = slide.shapes.add_textbox(left, top, width, height)
    tf = caixa.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = alinhar
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(tamanho)
    run.font.bold = bold
    run.font.color.rgb = cor
    run.font.name = fonte
    return caixa


def add_codigo(slide, linhas, left, top, width, height, tamanho=16, fg=COR_CODIGO_FG):
    bloco = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bloco.fill.solid()
    bloco.fill.fore_color.rgb = COR_CODIGO_BG
    bloco.line.color.rgb = COR_ACENTO
    bloco.line.width = Pt(1)
    bloco.adjustments[0] = 0.05

    caixa = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15),
        width - Inches(0.4), height - Inches(0.3),
    )
    tf = caixa.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    for i, linha in enumerate(linhas):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        cor_linha = fg
        if isinstance(linha, tuple):
            linha, cor_linha = linha
        run = p.add_run()
        run.text = linha if linha else " "
        run.font.name = "Consolas"
        run.font.size = Pt(tamanho)
        run.font.color.rgb = cor_linha


def add_bullet(slide, itens, left, top, width, height, tamanho=18):
    caixa = slide.shapes.add_textbox(left, top, width, height)
    tf = caixa.text_frame
    tf.word_wrap = True
    for i, item in enumerate(itens):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(tamanho)
        run.font.color.rgb = COR_TEXTO
        run.font.name = "Calibri"


def add_rodape(slide, numero, total):
    add_texto(
        slide,
        f"Versionamento com Git  •  Slide {numero}/{total}",
        Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3),
        tamanho=10, cor=RGBColor(0x88, 0x88, 0x99), alinhar=PP_ALIGN.LEFT,
    )


def slide_capa(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pintar_fundo(slide)

    add_texto(
        slide, "🧮",
        Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.5),
        tamanho=90, alinhar=PP_ALIGN.CENTER,
    )
    add_texto(
        slide, "Versionamento de Código com Git",
        Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.2),
        tamanho=48, bold=True, cor=COR_TITULO, alinhar=PP_ALIGN.CENTER,
    )
    add_texto(
        slide, "Um bug, uma branch, uma correção.",
        Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.8),
        tamanho=24, cor=COR_ACENTO, alinhar=PP_ALIGN.CENTER,
    )
    add_texto(
        slide, "Demo prática: calculadora em Python",
        Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.6),
        tamanho=16, cor=COR_TEXTO, alinhar=PP_ALIGN.CENTER,
    )
    add_faixa_acento(slide, topo=False)


def slide_conteudo(prs, numero, total, titulo, subtitulo=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pintar_fundo(slide)
    add_faixa_acento(slide, topo=True)
    add_texto(
        slide, titulo,
        Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.8),
        tamanho=32, bold=True, cor=COR_TITULO,
    )
    if subtitulo:
        add_texto(
            slide, subtitulo,
            Inches(0.6), Inches(1.15), Inches(12.1), Inches(0.5),
            tamanho=18, cor=COR_ACENTO,
        )
    add_rodape(slide, numero, total)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    TOTAL = 13

    # 1. Capa
    slide_capa(prs, TOTAL)

    # 2. O que é versionamento
    s = slide_conteudo(prs, 2, TOTAL, "O que é Versionamento?", "Controle de histórico do código")
    add_bullet(s, [
        "Registra TODA mudança feita no código ao longo do tempo",
        "Permite voltar pra qualquer versão anterior (ctrl+z turbinado)",
        "Vários devs trabalhando no mesmo projeto sem pisar no pé do outro",
        "Mostra QUEM mudou O QUE e QUANDO",
        "É a ferramenta mais importante do dia-a-dia do desenvolvedor",
    ], Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5), tamanho=20)

    # 3. Por que Git
    s = slide_conteudo(prs, 3, TOTAL, "Por que Git?", "O padrão da indústria")
    add_bullet(s, [
        "✅ Gratuito, open-source e funciona offline",
        "✅ Usado por 95%+ das empresas de tecnologia",
        "✅ Integra com GitHub, GitLab, Bitbucket…",
        "✅ Branches (ramificações) super baratas e rápidas",
        "✅ Histórico imutável — segurança e auditoria",
    ], Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5), tamanho=20)

    # 4. Conceitos fundamentais
    s = slide_conteudo(prs, 4, TOTAL, "Conceitos Fundamentais")
    add_bullet(s, [
        "Repositório: pasta versionada pelo Git",
        "Commit: uma 'foto' das mudanças, com mensagem",
        "Branch: uma linha de desenvolvimento paralela",
        "Merge: juntar uma branch com outra",
        "Pull Request (PR): pedido de revisão antes do merge",
        "Remote: cópia do repositório no servidor (GitHub)",
    ], Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), tamanho=20)

    # 5. Comandos essenciais
    s = slide_conteudo(prs, 5, TOTAL, "Comandos Essenciais", "O mínimo do mínimo")
    add_codigo(s, [
        "git clone <url>              # baixa um repositório",
        "git status                   # o que mudou?",
        "git checkout -b nova-branch  # cria e entra na branch",
        "git add arquivo.py           # prepara pro commit",
        "git commit -m \"msg\"          # registra o commit",
        "git push origin nova-branch  # envia pro GitHub",
        "git pull                     # baixa atualizações",
        "git log --oneline            # histórico resumido",
    ], Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5), tamanho=18)

    # 6. DEMO intro
    s = slide_conteudo(prs, 6, TOTAL, "🧪 Demo Prática", "Calculadora React com um bug simples")
    add_texto(s,
        "Quando a página carrega, o logo no topo aparece QUEBRADO.",
        Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.6), tamanho=20,
    )
    add_texto(s,
        "…o arquivo que o HTML tenta carregar não existe (404).",
        Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.6), tamanho=22, bold=True, cor=COR_VERMELHO,
    )
    add_bullet(s, [
        "Repositório: github.com/huyhian9-dotcom/git-demo-calculadora",
        "Pré-requisitos: Node.js 18+ e Git instalados",
        "Você vai investigar pelo DevTools (F12) → aba Network",
        "Correção: um único caractere!  Sério.",
        "Ao final, cada um vai ter enviado um Pull Request 🎯",
    ], Inches(0.8), Inches(3.6), Inches(11.7), Inches(3.0), tamanho=19)

    # 7. Passo 1 — Clone e rodar
    s = slide_conteudo(prs, 7, TOTAL, "Passo 1 — Clonar e rodar a aplicação")
    add_codigo(s, [
        "git clone https://github.com/huyhian9-dotcom/git-demo-calculadora.git",
        "cd git-demo-calculadora",
        "",
        "npm install      # instala React, Vite e afins",
        "npm run dev      # abre em http://localhost:5173",
    ], Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.6), tamanho=17)
    add_texto(s,
        "A calculadora abre no navegador — o logo no topo aparece QUEBRADO.",
        Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.6), tamanho=20, cor=COR_ACENTO,
    )

    # 8. NOVO — Passo 2: Abrir DevTools e identificar o bug
    s = slide_conteudo(prs, 8, TOTAL, "Passo 2 — Abrir o DevTools",
                       "Como devs, investigamos bugs pelas ferramentas do navegador")
    add_bullet(s, [
        "No navegador, pressione  F12  (ou Ctrl+Shift+I)",
        "Vá até a aba  Network",
        "Dê um refresh na página (F5)",
    ], Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.0), tamanho=19)
    add_texto(s,
        "Você verá uma requisição em VERMELHO:",
        Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.5), tamanho=18, bold=True, cor=COR_ACENTO,
    )
    add_codigo(s, [
        ("GET   /logoo.svg     404 Not Found  ❌", COR_VERMELHO),
    ], Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.0), tamanho=17)
    add_texto(s,
        "Percebeu o  \"logoo\"  com dois 'o'?  Esse é o bug.",
        Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.5), tamanho=16, cor=COR_TEXTO,
    )

    # 9. Passo 3 — Criar branch
    s = slide_conteudo(prs, 9, TOTAL, "Passo 3 — Criar uma branch de correção",
                       "Regra de ouro: nunca mexa direto na main")
    add_codigo(s, [
        "git checkout -b fix/logo-path",
    ], Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.1), tamanho=20)
    add_texto(s,
        "O que aconteceu?",
        Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.5), tamanho=20, bold=True, cor=COR_ACENTO,
    )
    add_bullet(s, [
        "Criou uma nova branch chamada fix/logo-path",
        "Trocou automaticamente pra ela (o -b faz isso)",
        "Agora qualquer mudança fica isolada, sem afetar a main",
    ], Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.5), tamanho=18)

    # 10. Passo 4 — Corrigir o caminho do logo
    s = slide_conteudo(prs, 10, TOTAL, "Passo 4 — Corrigir o caminho do logo",
                       "Arquivo: src/App.jsx")
    add_texto(s, "ANTES (com bug):",
        Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.4), tamanho=16, bold=True, cor=COR_VERMELHO,
    )
    add_codigo(s, [
        ('<img src="/logoo.svg"', COR_VERMELHO),
        ('     alt="Git Demo" />', COR_CODIGO_FG),
    ], Inches(0.8), Inches(2.4), Inches(5.8), Inches(1.5), tamanho=15)

    add_texto(s, "DEPOIS (corrigido):",
        Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.4), tamanho=16, bold=True, cor=COR_VERDE,
    )
    add_codigo(s, [
        ('<img src="/logo.svg"', COR_VERDE),
        ('     alt="Git Demo" />', COR_CODIGO_FG),
    ], Inches(7.0), Inches(2.4), Inches(5.8), Inches(1.5), tamanho=15)

    add_texto(s,
        "Removeu UM caractere.  Salve e dê refresh no navegador:",
        Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.5), tamanho=18, cor=COR_ACENTO,
    )
    add_codigo(s, [
        ("GET   /logo.svg     200 OK  ✅   (e o logo aparece!)", COR_VERDE),
    ], Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.0), tamanho=17)

    # 11. Passo 5 — Commit
    s = slide_conteudo(prs, 11, TOTAL, "Passo 5 — Registrar a correção (commit)")
    add_codigo(s, [
        "git status",
        "git add src/App.jsx",
        "git commit -m \"fix: corrige caminho do logo que retornava 404\"",
    ], Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.2), tamanho=16)
    add_texto(s,
        "Dica de mensagem de commit:",
        Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.5), tamanho=18, bold=True, cor=COR_ACENTO,
    )
    add_bullet(s, [
        "fix: corrige bug que estava quebrando X",
        "feat: adiciona funcionalidade Y",
        "docs: atualiza README",
        "refactor: reorganiza estrutura do arquivo Z",
    ], Inches(0.8), Inches(5.2), Inches(11.7), Inches(2.0), tamanho=16)

    # 12. Passo 6 — Push e PR
    s = slide_conteudo(prs, 12, TOTAL, "Passo 6 — Push + Pull Request")
    add_codigo(s, [
        "git push origin fix/logo-path",
    ], Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.1), tamanho=20)
    add_texto(s,
        "Depois, no GitHub:",
        Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.5), tamanho=18, bold=True, cor=COR_ACENTO,
    )
    add_bullet(s, [
        "Aparecerá o botão 'Compare & pull request' — clique nele",
        "Escreva um título claro e descreva o que você corrigiu",
        "Clique em 'Create pull request'",
        "Peça pro colega revisar — é assim que funciona no mundo real",
        "Após aprovação: 'Merge pull request' leva a correção pra main 🎉",
    ], Inches(0.8), Inches(4.0), Inches(11.7), Inches(3.0), tamanho=18)

    # 13. Encerramento
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pintar_fundo(slide)
    add_faixa_acento(slide, topo=True)
    add_faixa_acento(slide, topo=False)
    add_texto(slide, "🎯",
        Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.5),
        tamanho=80, alinhar=PP_ALIGN.CENTER,
    )
    add_texto(slide, "Você aprendeu o fluxo completo!",
        Inches(0.5), Inches(2.8), Inches(12.3), Inches(0.9),
        tamanho=40, bold=True, cor=COR_TITULO, alinhar=PP_ALIGN.CENTER,
    )
    add_texto(slide,
        "clone  →  branch  →  fix  →  commit  →  push  →  PR  →  merge",
        Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.7),
        tamanho=22, cor=COR_ACENTO, alinhar=PP_ALIGN.CENTER, fonte="Consolas",
    )
    add_texto(slide, "Perguntas?",
        Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.8),
        tamanho=30, bold=True, cor=COR_TEXTO, alinhar=PP_ALIGN.CENTER,
    )
    add_texto(slide, "github.com/huyhian9-dotcom/git-demo-calculadora",
        Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
        tamanho=14, cor=RGBColor(0x88, 0x88, 0x99), alinhar=PP_ALIGN.CENTER, fonte="Consolas",
    )

    saida = Path(__file__).parent / "apresentacao.pptx"
    prs.save(saida)
    print(f"[OK] Slides gerados: {saida}")


if __name__ == "__main__":
    build()
